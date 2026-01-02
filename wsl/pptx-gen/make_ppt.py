import argparse
import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches


def add_bullets(text_frame, bullets):
    text_frame.clear()
    if not bullets:
        return
    for i, item in enumerate(bullets):
        p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        p.text = str(item)
        p.level = 0


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title or ''
    if subtitle:
        slide.placeholders[1].text = subtitle


def add_content_slide(prs, spec):
    title = spec.get('title', '')
    bullets = spec.get('bullets', [])
    image = spec.get('image')
    notes = spec.get('notes')

    if image:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = title
        left = Inches(0.7)
        top = Inches(1.6)
        width = Inches(5.8)
        height = Inches(4.7)
        text_frame = slide.shapes.add_textbox(left, top, width, height).text_frame
        add_bullets(text_frame, bullets)
        img_path = Path(image)
        slide.shapes.add_picture(str(img_path), Inches(6.8), Inches(1.6), width=Inches(5.5))
    else:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        body = slide.shapes.placeholders[1].text_frame
        add_bullets(body, bullets)

    if notes:
        slide.notes_slide.notes_text_frame.text = str(notes)


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument('--input', required=True, help='Path to JSON spec')
    parser.add_argument('--output', required=True, help='Path to output .pptx')
    args = parser.parse_args()

    spec_path = Path(args.input)
    out_path = Path(args.output)

    with spec_path.open('r', encoding='utf-8') as f:
        spec = json.load(f)

    prs = Presentation()
    title = spec.get('title', 'Untitled Deck')
    subtitle = spec.get('subtitle')
    add_title_slide(prs, title, subtitle)

    for slide_spec in spec.get('slides', []):
        add_content_slide(prs, slide_spec)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)


if __name__ == '__main__':
    main()
